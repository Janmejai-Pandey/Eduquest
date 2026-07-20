import os
import pdfplumber
from pptx import Presentation
from docx import Document
from ocr import ocr_pdf, OCR_AVAILABLE

import json
import pickle
import re

import numpy as np
import faiss
from rank_bm25 import BM25Okapi
import torch
from sentence_transformers import SentenceTransformer

import llm_config as config


# Config
OCR_FALLBACK_MIN_CHARS = 30     # if page has fewer chars, try OCR


def parse_path_metadata(full_path: str) -> dict:
    """Extract branch/sem/subject/category from folder structure."""
    norm = os.path.normpath(full_path)
    parts = norm.split(os.sep)

    metadata = {
        "branch":   "",
        "semester": "",
        "subject":  "",
        "category": "",
        "rel_path": "",
    }

    try:
        sm_idx = parts.index("study_material")
        after = parts[sm_idx + 1:]
        metadata["rel_path"] = os.sep.join(after)

        if len(after) >= 1: metadata["branch"]   = after[0]
        if len(after) >= 2: metadata["semester"] = after[1]
        if len(after) >= 3: metadata["subject"]  = after[2]
        if len(after) >= 4: metadata["category"] = after[3]
    except ValueError:
        metadata["rel_path"] = os.path.basename(full_path)

    return metadata


def _build_record(path, text, location, file_type, meta):
    """Standardized record builder."""
    return {
        "source_file":     os.path.basename(path),
        "source_path":     os.path.normpath(path),
        "source_rel_path": meta["rel_path"],
        "branch":          meta["branch"],
        "semester":        meta["semester"],
        "subject":         meta["subject"],
        "category":        meta["category"],
        "file_type":       file_type,
        "location":        location,
        "text":            text,
    }


def _apply_meta_to_ocr_records(ocr_records: list, path: str, meta: dict) -> list:
    """Add full metadata to OCR records (which only have basic fields)."""
    for rec in ocr_records:
        rec["source_path"]     = os.path.normpath(path)
        rec["source_rel_path"] = meta["rel_path"]
        rec["branch"]          = meta["branch"]
        rec["semester"]        = meta["semester"]
        rec["subject"]         = meta["subject"]
        rec["category"]        = meta["category"]
    return ocr_records


# ═══════════════ PDF EXTRACTION ═══════════════

def extract_pdf(path: str, use_ocr_fallback: bool = True) -> list:
    """
    Extract text from PDF. Falls back to OCR for scanned/image pages.
    Shows a clean progress bar per PDF (no per-page clutter).
    """
    import sys

    records = []
    meta = parse_path_metadata(path)
    pages_needing_ocr = []
    total_pages = 0
    fname = os.path.basename(path)

    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()

                if text and len(text) >= OCR_FALLBACK_MIN_CHARS:
                    records.append(_build_record(
                        path, text, f"page {i}", "pdf", meta
                    ))
                else:
                    pages_needing_ocr.append(i)
    except Exception as e:
        print(f"[ERROR] Failed to read PDF {path}: {e}")
        if use_ocr_fallback and OCR_AVAILABLE:
            pages_needing_ocr = list(range(1, total_pages + 1)) if total_pages else None

    # ── OCR fallback with progress bar ──
    if use_ocr_fallback and OCR_AVAILABLE and pages_needing_ocr:
        pages_list = pages_needing_ocr if isinstance(pages_needing_ocr, list) else []
        total = len(pages_list) if pages_list else total_pages

        if total > 0:
            print(f"   🔍 OCR: {fname} — {total} scanned page(s)")

            for idx, page_num in enumerate(pages_list, start=1):
                # Show progress bar
                _print_progress(idx, total, prefix=f"      OCR")

                ocr_records = ocr_pdf(path, first_page=page_num, last_page=page_num, silent=True)
                records.extend(_apply_meta_to_ocr_records(ocr_records, path, meta))

            # Newline after progress completes
            print()

    return records


def _print_progress(current: int, total: int, prefix: str = "Progress", bar_length: int = 30):
    """Print/update a single-line progress bar."""
    import sys
    filled = int(bar_length * current // total)
    bar = '█' * filled + '░' * (bar_length - filled)
    percent = 100 * current // total
    sys.stdout.write(f"\r{prefix} [{bar}] {percent}% ({current}/{total})")
    sys.stdout.flush()

# ═══════════════ PPTX EXTRACTION ═══════════════

def extract_pptx(path: str) -> list:
    """Extract text from PPTX (including tables)."""
    records = []
    meta = parse_path_metadata(path)

    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

            # Also grab speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[NOTES] {notes}")

            text = "\n".join(parts).strip()
            if text:
                records.append(_build_record(
                    path, text, f"slide {i}", "pptx", meta
                ))
    except Exception as e:
        print(f"[ERROR] Failed to read PPTX {path}: {e}")
    return records


# ═══════════════ DOCX EXTRACTION ═══════════════

def extract_docx(path: str) -> list:
    """Extract text from a Word document (.docx)."""
    records = []
    meta = parse_path_metadata(path)

    try:
        doc = Document(path)
        parts = []

        # Paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        # Split by page-ish (Word doesn't have hard pages, so chunk by size)
        full_text = "\n".join(parts).strip()

        if full_text:
            # Treat entire doc as one "location" or split by ~500 word chunks
            words = full_text.split()
            words_per_page = 500

            if len(words) <= words_per_page:
                records.append(_build_record(
                    path, full_text, "document", "docx", meta
                ))
            else:
                # Split into pseudo-pages
                for i in range(0, len(words), words_per_page):
                    section_text = " ".join(words[i:i + words_per_page])
                    page_num = (i // words_per_page) + 1
                    records.append(_build_record(
                        path, section_text, f"section {page_num}", "docx", meta
                    ))
    except Exception as e:
        print(f"[ERROR] Failed to read DOCX {path}: {e}")

    return records


# ═══════════════ FOLDER SCAN ═══════════════

def extract_folder(folder: str) -> list:
    """Extract all PDFs, PPTXs, and DOCXs recursively."""
    all_records = []
    counts = {"pdf": 0, "pptx": 0, "docx": 0}

    for root, _, files in os.walk(folder):
        for fname in files:
            path = os.path.join(root, fname)
            lower = fname.lower()

            if lower.endswith(".pdf"):
                print(f"📄 Extracting PDF : {fname}")
                all_records.extend(extract_pdf(path))
                counts["pdf"] += 1
            elif lower.endswith(".pptx"):
                print(f"📊 Extracting PPTX: {fname}")
                all_records.extend(extract_pptx(path))
                counts["pptx"] += 1
            elif lower.endswith(".docx"):
                print(f"📝 Extracting DOCX: {fname}")
                all_records.extend(extract_docx(path))
                counts["docx"] += 1

    print(f"\n✅ Extracted {len(all_records)} pages/slides/sections from:")
    print(f"   • {counts['pdf']} PDF(s)")
    print(f"   • {counts['pptx']} PPTX(s)")
    print(f"   • {counts['docx']} DOCX(s)")

    return all_records


# ═══════════════ CHUNKING ═══════════════

def chunk_text(text: str, max_words: int = 200, overlap: int = 40) -> list:
    """Split text into overlapping word chunks."""
    words = text.split()
    if len(words) <= max_words:
        return [text]

    chunks = []
    start = 0
    while start < len(words):
        end = start + max_words
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end >= len(words):
            break
        start = end - overlap
    return chunks


def chunk_records(records: list, max_words: int = 200, overlap: int = 40) -> list:
    """Convert page/slide/section records into chunk records."""
    chunked = []
    for rec in records:
        pieces = chunk_text(rec["text"], max_words, overlap)
        for j, piece in enumerate(pieces):
            chunked.append({
                "chunk_id":        len(chunked),
                "source_file":     rec.get("source_file",     ""),
                "source_path":     rec.get("source_path",     ""),
                "source_rel_path": rec.get("source_rel_path", ""),
                "branch":          rec.get("branch",          ""),
                "semester":        rec.get("semester",        ""),
                "subject":         rec.get("subject",         ""),
                "category":        rec.get("category",        ""),
                "file_type":       rec.get("file_type",       ""),
                "location":        rec.get("location",        ""),
                "chunk_index":     j,
                "text":            piece,
            })
    print(f"✅ Created {len(chunked)} chunks with full metadata")
    return chunked

# ─────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────
CURRENT_DIR  = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(CURRENT_DIR, "..", ".."))
INDEX_DIR    = os.path.join(PROJECT_ROOT, "index_store")


# ─────────────────────────────────────────────
# Lazy model loader
# ─────────────────────────────────────────────
_embed_model = None
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
print(f"🖥️  Using device: {DEVICE.upper()}")
if DEVICE == "cuda":
    print(f"   GPU: {torch.cuda.get_device_name(0)}")


def get_embed_model():
    global _embed_model
    if _embed_model is None:
        print(f"📥 Loading embedding model: {config.EMBED_MODEL}")
        _embed_model = SentenceTransformer(config.EMBED_MODEL, device=DEVICE)
        print(f"✅ Embedding model loaded (dim={_embed_model.get_sentence_embedding_dimension()})")
    return _embed_model


# ─────────────────────────────────────────────
# Tokenizer for BM25
# ─────────────────────────────────────────────
def tokenize(text: str) -> list:
    return re.findall(r"[a-z0-9]+", text.lower())


# ─────────────────────────────────────────────
# Build indexes
# ─────────────────────────────────────────────
def build_indexes(chunks: list):
    os.makedirs(INDEX_DIR, exist_ok=True)
    texts = [c["text"] for c in chunks]

    # BM25
    print("\n🔨 Building BM25 index...")
    tokenized = [tokenize(t) for t in texts]
    bm25 = BM25Okapi(tokenized)
    with open(os.path.join(INDEX_DIR, "bm25.pkl"), "wb") as f:
        pickle.dump({"bm25": bm25, "tokenized": tokenized}, f)
    print(f"✅ BM25 index built ({len(tokenized)} docs)")

    # Vector (BGE-M3 on GPU)
    print(f"\n🔨 Building vector index with {config.EMBED_MODEL} on {DEVICE.upper()}...")
    model = get_embed_model()

    # On GPU can use larger batch size
    batch_size = 64 if DEVICE == "cuda" else config.EMBED_BATCH_SIZE

    embeddings = model.encode(
        texts,
        batch_size            = batch_size,
        show_progress_bar     = True,
        convert_to_numpy      = True,
        normalize_embeddings  = config.EMBED_NORMALIZE,
    )

    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings.astype(np.float32))
    faiss.write_index(index, os.path.join(INDEX_DIR, "faiss.index"))
    print(f"✅ FAISS index built (dim={dim})")

    # Save chunks
    with open(os.path.join(INDEX_DIR, "chunks.json"), "w", encoding="utf-8") as f:
        json.dump(chunks, f, ensure_ascii=False)

    # Save config
    idx_config = {
        "embed_model": config.EMBED_MODEL,
        "embed_dim":   dim,
        "num_chunks":  len(chunks),
        "device":      DEVICE,
    }
    with open(os.path.join(INDEX_DIR, "config.json"), "w", encoding="utf-8") as f:
        json.dump(idx_config, f, indent=2)

    print(f"\n✅ All indexes saved to {INDEX_DIR}") 
    
# =================Main function to run the whole pipeline===================

if __name__ == "__main__":
    folder = "dataset/study_material"  # Change this to your folder path
    print(f"Extracting documents from folder: {folder}")
    records = extract_folder(folder)

    print("Chunking extracted text...")
    chunks = chunk_records(records)

    print("Building indexes...")
    build_indexes(chunks)

    print("All done! You can now query the indexes for retrieval.")
