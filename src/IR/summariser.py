"""
summariser.py
Lecture PDF / PPTX summariser.
Streams files directly from Google Drive (no local download).
"""

import os
import io
import json
import re
from typing import List, Dict, Optional, Tuple

import requests
import pdfplumber
from pptx import Presentation

from llm import get_answer


# ============== PATHS ==============

BASE_DATASET = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "dataset", "study_material", "CSE"
)

SEM_FOLDERS = {
    "3": os.path.join(BASE_DATASET, "3"),
    "4": os.path.join(BASE_DATASET, "4"),
}


# ============== PROMPTS ==============

SINGLE_LECTURE_PROMPT = """You are an expert academic summariser.
Summarise the following lecture content thoroughly and precisely.

Your summary MUST include:

## Summary
A clear, structured summary covering every major concept discussed.

## Key Formulae & Equations
List ALL mathematical formulae, equations, theorems, or formal
definitions found in the text. Use proper notation.
If none exist, write "No formulae found in this lecture."

## Important Points
Numbered list of every critical concept, definition, or fact
a student must know.

## Quick Memorisation Bullets
- 8-15 concise bullet points a student can revise in 5 minutes
  before an exam. Focus on what is most likely to be asked.

Be thorough - do NOT skip content. If the lecture is long,
cover everything.
"""

SERIES_LECTURE_PROMPT = """You are an expert academic summariser.
You are given the combined content of MULTIPLE lecture notes
(in order). Produce a UNIFIED summary across all of them.

Your summary MUST include:

## Combined Summary
A comprehensive summary spanning all the lectures provided.
Organise by topic / theme, not by individual lecture.

## All Formulae & Equations
Collect ALL mathematical formulae, equations, theorems, or
formal definitions across all lectures. Group by topic.
If none exist, write "No formulae found."

## Important Points (across all lectures)
Numbered list of every critical concept, definition, or fact.

## Lecture Flow & Connections
Briefly describe how the topics connect from one lecture to the next.

## Quick Memorisation Bullets
- 12-20 concise bullet points covering the entire lecture series.
  Focus on what is most likely to be asked in exams.

Be thorough - cover ALL lectures provided.
"""


# ============== LOAD & FILTER ==============

def load_gdrive_links(semester: str) -> List[Dict]:
    """Load gdrive_links.json for the given semester."""
    sem_folder = SEM_FOLDERS.get(semester)
    if not sem_folder or not os.path.isdir(sem_folder):
        raise FileNotFoundError(f"Semester folder not found: {sem_folder}")

    json_path = os.path.join(sem_folder, "gdrive_links.json")
    if not os.path.isfile(json_path):
        raise FileNotFoundError(f"gdrive_links.json not found at: {json_path}")

    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_subjects(links: List[Dict]) -> List[str]:
    """Return sorted unique subject names."""
    return sorted(set(link["subject"] for link in links))


def get_lecture_files(links: List[Dict], subject: str) -> List[Dict]:
    """Return only 'Lectures' for a subject, naturally sorted."""
    lectures = [
        link for link in links
        if link["subject"] == subject
        and link["category"] == "Lectures"
        and not link.get("is_folder", False)
    ]

    def natural_sort_key(item):
        name = item.get("item_name", "")
        numbers = re.findall(r"\d+", name)
        first_num = int(numbers[0]) if numbers else 999
        return (first_num, name.lower())

    lectures.sort(key=natural_sort_key)
    return lectures


# ============== STREAM FILE FROM GOOGLE DRIVE ==============

def stream_from_gdrive(file_id: str) -> Optional[io.BytesIO]:
    """
    Download a Google Drive file into an in-memory buffer.
    Returns BytesIO object or None on failure.
    """
    session = requests.Session()
    session.headers.update({
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36"
        )
    })

    base_url = "https://drive.google.com/uc?export=download"

    try:
        # First request
        response = session.get(
            base_url, params={"id": file_id}, stream=True, timeout=30
        )

        # Handle Google's "confirm download" page for large files
        token = None
        for k, v in response.cookies.items():
            if k.startswith("download_warning"):
                token = v
                break

        if not token:
            # Sometimes token is in HTML body
            match = re.search(r'confirm=([0-9A-Za-z_\-]+)', response.text)
            if match:
                token = match.group(1)

        if token:
            response = session.get(
                base_url,
                params={"id": file_id, "confirm": token},
                stream=True,
                timeout=30,
            )

        # Verify we got a binary file, not an HTML error page
        content_type = response.headers.get("Content-Type", "")
        if "text/html" in content_type:
            print(f"   Warning: Got HTML response (file may be too large or restricted)")

        # Read into buffer
        buffer = io.BytesIO()
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                buffer.write(chunk)

        buffer.seek(0)
        return buffer

    except Exception as e:
        print(f"   Stream failed: {e}")
        return None


# ============== EXTRACT FROM BUFFER ==============

def extract_pdf_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    """Extract text page-wise from a PDF in memory."""
    records = []
    try:
        with pdfplumber.open(buffer) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                text = text.strip()
                if text:
                    records.append({
                        "source_file": source_name,
                        "file_type": "pdf",
                        "location": f"page {i}",
                        "text": text,
                    })
    except Exception as e:
        print(f"   PDF extraction failed: {e}")
    return records


def extract_pptx_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    """Extract text slide-wise from a PPTX in memory."""
    records = []
    try:
        prs = Presentation(buffer)
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            text = "\n".join(parts).strip()
            if text:
                records.append({
                    "source_file": source_name,
                    "file_type": "pptx",
                    "location": f"slide {i}",
                    "text": text,
                })
    except Exception as e:
        print(f"   PPTX extraction failed: {e}")
    return records


def extract_text_from_gdrive(lecture_info: Dict) -> Tuple[str, str]:
    """
    Stream file from Google Drive and extract its text.
    Returns (extracted_text, view_url).
    """
    file_id = lecture_info.get("file_id", "")
    item_name = lecture_info.get("item_name", "Unknown")
    view_url = lecture_info.get(
        "original_url",
        f"https://drive.google.com/file/d/{file_id}/view"
    )

    if not file_id:
        return "", view_url

    print(f"   Streaming from Google Drive...")
    buffer = stream_from_gdrive(file_id)
    if not buffer:
        return "", view_url

    lower_name = item_name.lower()
    records = []

    if lower_name.endswith(".pdf"):
        records = extract_pdf_from_buffer(buffer, item_name)
    elif lower_name.endswith(".pptx"):
        records = extract_pptx_from_buffer(buffer, item_name)
    else:
        # Try PDF first (some files have no extension in name)
        buffer.seek(0)
        records = extract_pdf_from_buffer(buffer, item_name)
        if not records:
            buffer.seek(0)
            records = extract_pptx_from_buffer(buffer, item_name)

    if not records:
        return "", view_url

    parts = []
    for rec in records:
        parts.append(f"--- {rec['location']} ---\n{rec['text']}")

    return "\n\n".join(parts), view_url


# ============== SUMMARISE ==============

def summarise_single_lecture(
    semester: str, lecture_info: Dict
) -> Tuple[str, bool, str]:
    """
    Summarise a single lecture file streamed from Google Drive.
    Returns (summary_text, success_flag, view_url).
    """
    item_name = lecture_info.get("item_name", "Unknown")
    print(f"\nProcessing: {item_name}")

    text, view_url = extract_text_from_gdrive(lecture_info)

    if not text.strip():
        return (
            f"Could not extract text from: {item_name}\n"
            f"   View the file directly here: {view_url}",
            False,
            view_url,
        )

    word_count = len(text.split())
    print(f"   Extracted {word_count} words")

    MAX_WORDS = 12000
    if word_count > MAX_WORDS:
        words = text.split()
        text = " ".join(words[:MAX_WORDS])
        print(f"   Truncated to {MAX_WORDS} words for LLM context")

    messages = [
        {"role": "system", "content": SINGLE_LECTURE_PROMPT},
        {
            "role": "user",
            "content": (
                f"Lecture: {item_name}\n"
                f"Subject: {lecture_info.get('subject', 'Unknown')}\n\n"
                f"Full Lecture Content:\n{text}"
            ),
        },
    ]

    print("   Generating summary...")
    summary = get_answer(messages)
    return summary, True, view_url


def summarise_lecture_series(
    semester: str,
    lecture_list: List[Dict],
    subject: str,
) -> Tuple[str, bool, List[Dict]]:
    """
    Summarise a series of lectures into one combined summary.
    Returns (summary_text, success_flag, list_of_lecture_links).
    """
    print(f"\nProcessing {len(lecture_list)} lectures for: {subject}")

    all_texts = []
    processed_links = []  # list of {"name": ..., "url": ...}
    processed = 0

    for lec in lecture_list:
        item_name = lec.get("item_name", "Unknown")
        print(f"\n   Lecture: {item_name}")

        text, view_url = extract_text_from_gdrive(lec)

        if not text.strip():
            print(f"   Skipping (no text extracted): {item_name}")
            continue

        all_texts.append(
            f"\n{'='*60}\n"
            f"LECTURE: {item_name}\n"
            f"{'='*60}\n"
            f"{text}"
        )
        processed_links.append({"name": item_name, "url": view_url})
        processed += 1
        print(f"   Extracted: {item_name}")

    if not all_texts:
        return "No lecture files could be extracted.", False, processed_links

    combined_text = "\n\n".join(all_texts)
    total_words = len(combined_text.split())
    print(f"\n   Total: {processed} lectures, {total_words} words")

    MAX_WORDS = 15000
    if total_words > MAX_WORDS:
        words = combined_text.split()
        combined_text = " ".join(words[:MAX_WORDS])
        print(f"   Truncated to {MAX_WORDS} words for LLM context")

    messages = [
        {"role": "system", "content": SERIES_LECTURE_PROMPT},
        {
            "role": "user",
            "content": (
                f"Subject: {subject}\n"
                f"Lectures included: "
                f"{', '.join(l['item_name'] for l in lecture_list)}\n\n"
                f"Combined Lecture Content:\n{combined_text}"
            ),
        },
    ]

    print("   Generating combined summary...")
    summary = get_answer(messages)
    return summary, True, processed_links


# ============== SAVE ==============

def save_summary(
    summary: str,
    semester: str,
    subject: str,
    label: str,
    links: Optional[List[Dict]] = None,
) -> str:
    """Save summary to a markdown file. Returns the saved path."""
    safe_subject = re.sub(r'[<>:"/\\|?*]', "_", subject).strip()
    out_dir = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "..", "..", "summaries",
        f"sem_{semester}",
        safe_subject,
    )
    os.makedirs(out_dir, exist_ok=True)

    clean_label = re.sub(r'[<>:"/\\|?*]', "_", label).strip()
    # Trim very long filenames
    if len(clean_label) > 100:
        clean_label = clean_label[:100]
    filename = f"summary_{clean_label}.md"
    filepath = os.path.join(out_dir, filename)

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(f"# Summary: {label}\n")
        f.write(f"**Subject**: {subject}\n")
        f.write(f"**Semester**: {semester}\n\n")

        if links:
            f.write("## Source Lectures (click to open)\n\n")
            for link in links:
                f.write(f"- [{link['name']}]({link['url']})\n")
            f.write("\n")

        f.write("---\n\n")
        f.write(summary)

    return filepath