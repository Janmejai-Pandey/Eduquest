"""
pyq_analyser.py
Previous Year Questions analyser.
Streams PYQ papers directly from Google Drive (no local downloads).
Handles subject aliases, lab vs theory, and content verification.
"""

import os
import io
import json
import re
from typing import List, Dict, Optional, Tuple

import requests
import pdfplumber
from pptx import Presentation

from llm_env import get_answer
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# ============== PATHS ==============

BASE_DATASET = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "dataset", "study_material", "CSE"
)

SEM_FOLDERS = {
    "3": os.path.join(BASE_DATASET, "3"),
    "4": os.path.join(BASE_DATASET, "4"),
}


# ============== SUBJECT ALIASES ==============
# Maps user-typed short forms to canonical subject names in gdrive_links.json
# All keys are lowercase for case-insensitive matching.

SUBJECT_ALIASES = {
    # Mathematical Foundations for AI and DS
    "maths":     "Mathematical Foundations for AI and DS",
    "math":      "Mathematical Foundations for AI and DS",
    "mathematics": "Mathematical Foundations for AI and DS",
    "mfaids":    "Mathematical Foundations for AI and DS",
    "mf":        "Mathematical Foundations for AI and DS",

    # OOP using Java Lab
    "oops":      "OOP using Java Lab",
    "oop":       "OOP using Java Lab",
    "java":      "OOP using Java Lab",
    "java lab":  "OOP using Java Lab",
    "oops lab":  "OOP using Java Lab",

    # Unix Programming Lab
    "unix":      "Unix Programming Lab",
    "unix lab":  "Unix Programming Lab",
    "up":        "Unix Programming Lab",

    # Database Management Systems (theory) + DBMS Lab
    "dbms":      "Database Management Systems",           # default → theory
    "dbms theory": "Database Management Systems",
    "dbms lab":  "Database Management Systems Lab",
    "database":  "Database Management Systems",
    "database management": "Database Management Systems",

    # Theory of Computation
    "toc":       "Theory of Computation",
    "theory of computation": "Theory of Computation",
    "automata":  "Theory of Computation",

    # Economics
    "eco":       "Economics",
    "econ":      "Economics",
    "economics": "Economics",

    # Data Structures (theory) + Lab
    "ds":        "Data Structures",                       # default → theory
    "ds theory": "Data Structures",
    "ds lab":    "Data Structures Lab",
    "data structures": "Data Structures",
    "dsa":       "Data Structures",

    # Sem 4 additions (some coverage)
    "dsco":      "Digital Systems and Comp. Organisation",
    "se":        "Software Engineering",
    "software":  "Software Engineering",
    "aiml":      "Artificial Intelligence and Machine Learning",
    "ai ml":     "Artificial Intelligence and Machine Learning",
    "ml":        "Artificial Intelligence and Machine Learning",
    "ai":        "Artificial Intelligence and Machine Learning",
    "daa":       "Design and Analysis of Algorithms",
    "algo":      "Design and Analysis of Algorithms",
    "evs":       "Environmental Studies Lab",
    "qmss":      "Quantitative Methods for Social Sc. (Elective)",
    "mad":       "Fundamentals of Mobile Appl. Development (Elective)",
    "mad lab":   "Fundamentals of Mobile Appl. Development Lab (Elective)",
    "aiml lab":  "Artificial Intelligence and Machine Learning Lab",
    "daa lab":   "Design and Analysis of Algorithms Lab",
    "cp":        "Competitive Programming Lab",
}


# Subjects that have BOTH a Lab and Theory version — ask user which
SUBJECTS_WITH_LAB_THEORY = {
    "dbms":  ("Database Management Systems",     "Database Management Systems Lab"),
    "ds":    ("Data Structures",                 "Data Structures Lab"),
    "aiml":  ("Artificial Intelligence and Machine Learning",
              "Artificial Intelligence and Machine Learning Lab"),
    "daa":   ("Design and Analysis of Algorithms",
              "Design and Analysis of Algorithms Lab"),
    "mad":   ("Fundamentals of Mobile Appl. Development (Elective)",
              "Fundamentals of Mobile Appl. Development Lab (Elective)"),
}


# ============== EXAM TYPE DETECTION ==============

def detect_exam_type(filename: str) -> Optional[str]:
    """Detect exam type from PYQ filename.
    Handles T1/T2/T3, M1/M2/M3 (midterm), and lab-specific patterns."""
    name = filename.upper()

    # ── Priority 1: T3 / M3 / End Term ──────────────────────────────────
    if re.search(r"(?:^|[\s_\-\.])T3(?:[\s_\-\.]|$)", name):
        return "T3"
    if re.search(r"(?:^|[\s_\-\.])M3(?:[\s_\-\.]|$)", name):
        return "T3"
    if "ENDTERM" in name or "END TERM" in name or "END-TERM" in name or "END_TERM" in name:
        return "T3"

    # ── Priority 2: T2 / M2 ─────────────────────────────────────────────
    if re.search(r"(?:^|[\s_\-\.])T2(?:[\s_\-\.]|$)", name):
        return "T2"
    if re.search(r"(?:^|[\s_\-\.])M2(?:[\s_\-\.]|$)", name):
        return "T2"

    # ── Priority 3: T1 / M1 ─────────────────────────────────────────────
    if re.search(r"(?:^|[\s_\-\.])T1(?:[\s_\-\.]|$)", name):
        return "T1"
    if re.search(r"(?:^|[\s_\-\.])M1(?:[\s_\-\.]|$)", name):
        return "T1"

    # ── Combined papers ─────────────────────────────────────────────────
    if "MIDTERM" in name and "ENDTERM" in name:
        return "T3"

    # ── Lab-specific patterns ────────────────────────────────────────────
    if "LAB TEST-2" in name or "LABTEST2" in name or "LAB TEST 2" in name or "LAB-TEST-2" in name:
        return "T2"
    if "LAB TEST-1" in name or "LABTEST1" in name or "LAB TEST 1" in name or "LAB-TEST-1" in name:
        return "T1"

    if re.search(r"EVAL\s*#?\s*0?2|EVAL[-_ ]?2|EVAL[-_ ]?II|LAB\s*EVAL\s*2", name):
        return "T2"
    if re.search(r"EVAL\s*#?\s*0?1|EVAL[-_ ]?1|EVAL[-_ ]?I(?!I)|LAB\s*EVAL\s*1", name):
        return "T1"

    if "MIDTERM" in name:
        return "T2"

    return None

def get_pyq_files_for_exam(pyq_files: List[Dict], exam: str) -> List[Dict]:
    """Filter PYQ files by exam type (T1/T2/T3)."""
    matched = []
    for f in pyq_files:
        name = f.get("item_name", "")
        exam_type = detect_exam_type(name)
        if exam_type == exam:
            matched.append(f)
        elif exam in ("T1", "T2") and re.search(rf"T1[_\s]*T2|T1\+T2", name.upper()):
            matched.append(f)
    return matched


# ============== LOAD & FILTER ==============

def load_gdrive_links(semester: str) -> List[Dict]:
    """Load gdrive_links.json for the given semester."""
    sem_folder = SEM_FOLDERS.get(semester)
    if not sem_folder or not os.path.isdir(sem_folder):
        raise FileNotFoundError(f"Semester folder not found: {sem_folder}")
    json_path = os.path.join(sem_folder, "gdrive_links.json")
    if not os.path.isfile(json_path):
        raise FileNotFoundError(f"gdrive_links.json not found at: {json_path}")
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_subjects_with_pyqs(links: List[Dict]) -> List[str]:
    """Return subjects that have at least one PYQ file."""
    subjects = set()
    for link in links:
        if link.get("category") == "PYQs" and not link.get("is_folder", False):
            subjects.add(link["subject"])
    return sorted(subjects)


def resolve_subject(
    user_input: str,
    available_subjects: List[str],
) -> Optional[Tuple[str, Optional[str]]]:
    """
    Resolve user's input to a canonical subject name.

    Returns:
        (subject_name, needs_lab_theory_choice_key) if resolvable.
        (subject_name, None) if direct match.
        None if not resolvable.

    'needs_lab_theory_choice_key' is the alias key (e.g., 'dbms') if
    caller should ask user "Lab or Theory?".
    """
    q = user_input.strip().lower()

    # Empty input
    if not q:
        return None

    # Try exact alias match first
    if q in SUBJECT_ALIASES:
        canonical = SUBJECT_ALIASES[q]

        # Is this an ambiguous alias (lab vs theory)?
        if q in SUBJECTS_WITH_LAB_THEORY:
            return (canonical, q)  # signal that user needs to be asked

        if canonical in available_subjects:
            return (canonical, None)
        # Alias mapped, but subject not in this sem — try fuzzy
        for subj in available_subjects:
            if canonical.lower() in subj.lower() or subj.lower() in canonical.lower():
                return (subj, None)

    # Try exact match against available subjects (case-insensitive)
    for subj in available_subjects:
        if subj.lower() == q:
            return (subj, None)

    # Try substring match
    matches = [s for s in available_subjects if q in s.lower()]
    if len(matches) == 1:
        return (matches[0], None)

    # Fuzzy: any word in query matches subject
    query_words = set(q.split())
    scored = []
    for subj in available_subjects:
        subj_words = set(subj.lower().split())
        overlap = len(query_words & subj_words)
        if overlap > 0:
            scored.append((overlap, subj))
    scored.sort(reverse=True)
    if scored:
        return (scored[0][1], None)

    return None


def get_pyq_files(links: List[Dict], subject: str) -> List[Dict]:
    """Return PYQ files for a subject."""
    return [
        link for link in links
        if link.get("subject") == subject
        and link.get("category") == "PYQs"
        and not link.get("is_folder", False)
    ]


def get_lecture_context(links: List[Dict], subject: str, max_chars: int = 4000) -> str:
    """Gather brief syllabus overview from lecture filenames."""
    lectures = [
        link.get("item_name", "") for link in links
        if link.get("subject") == subject and link.get("category") == "Lectures"
    ]
    if not lectures:
        return "(no lecture syllabus available)"
    summary = "Lecture topics covered in this subject:\n"
    for lec in lectures:
        summary += f"- {lec}\n"
        if len(summary) > max_chars:
            break
    return summary


# ============== STREAM FROM GDRIVE ==============

def stream_from_gdrive(file_id: str) -> Optional[io.BytesIO]:
    session = requests.Session()
    session.headers.update({
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36"
        )
    })
    base_url = "https://drive.google.com/uc?export=download"
    try:
        response = session.get(base_url, params={"id": file_id}, stream=True, timeout=30)
        token = None
        for k, v in response.cookies.items():
            if k.startswith("download_warning"):
                token = v
                break
        if not token:
            match = re.search(r'confirm=([0-9A-Za-z_\-]+)', response.text)
            if match:
                token = match.group(1)
        if token:
            response = session.get(
                base_url, params={"id": file_id, "confirm": token},
                stream=True, timeout=30,
            )
        buffer = io.BytesIO()
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                buffer.write(chunk)
        buffer.seek(0)
        return buffer
    except Exception as e:
        print(f"   Stream failed: {e}")
        return None


def extract_pdf_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    records = []
    try:
        with pdfplumber.open(buffer) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    records.append({
                        "source_file": source_name,
                        "location": f"page {i}",
                        "text": text,
                    })
    except Exception as e:
        print(f"   PDF extraction failed: {e}")
    return records


def extract_pptx_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    records = []
    try:
        prs = Presentation(buffer)
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            text = "\n".join(parts).strip()
            if text:
                records.append({
                    "source_file": source_name,
                    "location": f"slide {i}",
                    "text": text,
                })
    except Exception as e:
        print(f"   PPTX extraction failed: {e}")
    return records

def extract_docx_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    """Extract text from a .docx file in memory."""
    records = []
    if not DOCX_AVAILABLE:
        print(f"   ⚠ python-docx not installed — cannot read {source_name}")
        return records

    try:
        doc = Document(buffer)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Also read tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            records.append({
                "source_file": source_name,
                "location": "document",
                "text": text,
            })
    except Exception as e:
        print(f"   DOCX extraction failed: {e}")
    return records

def extract_text_from_pyq(file_info: Dict) -> str:
    """Stream file from Drive and extract text. Handles PDF, PPTX, DOCX,
    and files without extensions."""
    file_id = file_info.get("file_id", "")
    item_name = file_info.get("item_name", "Unknown")
    if not file_id:
        return ""

    print(f"   Streaming: {item_name}")
    buffer = stream_from_gdrive(file_id)
    if not buffer:
        print(f"      ✗ Could not stream from Drive")
        return ""

    # Check if we accidentally got an HTML error page from Drive
    buffer.seek(0)
    first_bytes = buffer.read(500)
    buffer.seek(0)

    if first_bytes.startswith(b"<!DOCTYPE") or first_bytes.startswith(b"<html"):
        print(f"      ✗ Drive returned HTML error page (file may be too big or restricted)")
        return ""

    lower_name = item_name.lower()
    records = []

    # ── Try by extension first ──
    if lower_name.endswith(".pdf"):
        records = extract_pdf_from_buffer(buffer, item_name)
    elif lower_name.endswith(".pptx"):
        records = extract_pptx_from_buffer(buffer, item_name)
    elif lower_name.endswith(".docx"):
        records = extract_docx_from_buffer(buffer, item_name)
    else:
        # ── No extension or unknown — try all in sequence ──
        # PDF (most common)
        records = extract_pdf_from_buffer(buffer, item_name)

        # PPTX
        if not records:
            buffer.seek(0)
            records = extract_pptx_from_buffer(buffer, item_name)

        # DOCX
        if not records:
            buffer.seek(0)
            records = extract_docx_from_buffer(buffer, item_name)

    if not records:
        print(f"      ✗ Could not extract text from {item_name} "
              f"(file may be corrupt or unsupported format)")
        return ""

    return "\n\n".join(f"--- {r['location']} ---\n{r['text']}" for r in records)

# ============== SUBJECT VERIFICATION ==============

# Keywords that identify each subject when found inside a paper's text
SUBJECT_KEYWORDS = {
    "Mathematical Foundations for AI and DS": [
        "mathematical foundations", "mfaids", "mf ai", "math foundations",
        "probability", "random variable", "linear algebra", "vector space",
        "regression", "hypothesis testing",
    ],
    "OOP using Java Lab": [
        "java", "object oriented", "class", "inheritance", "polymorphism",
        "encapsulation", "constructor", "abstract class", "interface",
    ],
    "Unix Programming Lab": [
        "unix", "shell script", "shell scripting", "bash", "sed", "awk",
        "grep", "chmod", "chown", "linux command",
    ],
    "Database Management Systems": [
        "dbms", "database management", "sql", "normalization", "er diagram",
        "relational", "primary key", "foreign key", "join", "transaction",
    ],
    "Database Management Systems Lab": [
        "dbms", "sql", "create table", "insert into", "select", "join",
        "database", "relational", "trigger", "stored procedure",
    ],
    "Theory of Computation": [
        "theory of computation", "toc", "finite automata", "dfa", "nfa",
        "regular expression", "turing machine", "pushdown automata",
        "context free", "pumping lemma", "grammar",
    ],
    "Economics": [
        "economics", "demand", "supply", "elasticity", "gdp", "inflation",
        "market", "monopoly", "consumer", "producer",
    ],
    "Data Structures": [
        "data structure", "linked list", "stack", "queue", "tree",
        "binary tree", "heap", "graph", "hashing", "sorting", "algorithm",
    ],
    "Data Structures Lab": [
        "data structure", "linked list", "stack", "queue", "tree",
        "sorting", "search", "algorithm", "implement",
    ],
    "Digital Systems and Comp. Organisation": [
        "digital", "boolean", "karnaugh", "flip flop", "counter",
        "register", "computer organisation", "instruction",
    ],
    "Software Engineering": [
        "software engineering", "sdlc", "agile", "waterfall", "requirement",
        "uml", "use case", "sprint",
    ],
    "Artificial Intelligence and Machine Learning": [
        "artificial intelligence", "machine learning", "neural network",
        "supervised", "unsupervised", "classification", "regression",
        "clustering", "decision tree",
    ],
    "Design and Analysis of Algorithms": [
        "algorithm", "time complexity", "space complexity", "asymptotic",
        "dynamic programming", "greedy", "divide and conquer", "np hard",
    ],
    "Environmental Studies Lab": [
        "environmental", "pollution", "ecosystem", "biodiversity",
        "climate", "sustainability",
    ],
    "Quantitative Methods for Social Sc. (Elective)": [
        "quantitative", "statistics", "hypothesis", "regression", "sampling",
        "distribution",
    ],
    "Fundamentals of Mobile Appl. Development (Elective)": [
        "android", "mobile", "activity", "intent", "kotlin", "java",
        "layout", "fragment", "manifest",
    ],
    "Fundamentals of Mobile Appl. Development Lab (Elective)": [
        "android", "mobile", "activity", "intent", "kotlin", "xml layout",
    ],
    "Competitive Programming Lab": [
        "competitive programming", "algorithm", "time complexity",
        "problem solving",
    ],
    "Artificial Intelligence and Machine Learning Lab": [
        "machine learning", "python", "numpy", "pandas", "scikit",
        "classifier", "training",
    ],
}


def verify_paper_belongs_to_subject(
    paper_text: str, subject: str, threshold: int = 1
) -> bool:
    """
    Check if the extracted paper text actually belongs to this subject.
    Uses keyword matching against SUBJECT_KEYWORDS.
    Returns True if at least `threshold` keywords are found.
    """
    if not paper_text or not paper_text.strip():
        return False

    keywords = SUBJECT_KEYWORDS.get(subject, [])
    if not keywords:
        # No keywords defined for this subject — accept by default
        return True

    text_lower = paper_text.lower()
    matches = sum(1 for kw in keywords if kw.lower() in text_lower)

    return matches >= threshold


def gather_pyq_content(
    pyq_files: List[Dict], subject: str
) -> Tuple[str, List[str], List[str]]:
    """
    Extract text from all PYQ files and VERIFY they belong to the subject.
    Returns (combined_text, list_of_included_filenames, list_of_rejected_filenames).
    """
    all_texts = []
    included = []
    rejected = []

    for f in pyq_files:
        text = extract_text_from_pyq(f)
        fname = f.get("item_name", "Unknown")

        if not text.strip():
            print(f"   ⚠ Skipped (no text extracted): {fname}")
            rejected.append(fname)
            continue

        # Verify subject match
        if not verify_paper_belongs_to_subject(text, subject):
            print(f"   ⚠ Rejected (content doesn't match subject '{subject}'): {fname}")
            rejected.append(fname)
            continue

        all_texts.append(
            f"\n{'='*60}\nPAPER: {fname}\n{'='*60}\n{text}"
        )
        included.append(fname)
        print(f"   ✓ Included: {fname}")

    return "\n\n".join(all_texts), included, rejected


# ============== ANALYSIS PROMPTS ==============

def build_analysis_prompt(subject: str, exam: str, num_papers: int) -> str:
    exam_context = {
        "T1": "T1 (Test 1) - first ~30% of syllabus",
        "T2": "T2 (Test 2) - middle ~30% of syllabus (weeks 6-10)",
        "T3": "T3 (End Semester) - full syllabus, mix of T1+T2+new content",
    }.get(exam, exam)

    return f"""You are an expert academic analyst specialising in exam paper analysis.

You will be given the FULL TEXT of {num_papers} previous year question paper(s)
for the subject: **{subject}**
Exam type: **{exam_context}**

Your task: Perform a THOROUGH analysis and produce a structured JSON output
with the following schema (respond with ONLY valid JSON, no extra text):

{{
  "papers_analyzed": <int>,
  "years_covered": "<string like '2020-2025' or 'various'>",

  "topic_frequency": [
    {{
      "topic": "<topic name>",
      "count": <int - how many times asked>,
      "typical_marks": <int - average marks when asked>,
      "priority": "HIGH | MEDIUM | LOW"
    }}
  ],

  "question_type_distribution": {{
    "long_answer_percent": <int 0-100>,
    "short_answer_percent": <int 0-100>,
    "mcq_percent": <int 0-100>,
    "numerical_percent": <int 0-100>
  }},

  "marks_pattern": {{
    "total_marks_typical": <int>,
    "most_common_question_marks": <int>,
    "notes": "<brief string about marks distribution pattern>"
  }},

  "recurring_questions": [
    {{
      "question_pattern": "<the recurring question or pattern>",
      "frequency": <int>,
      "topic": "<topic it belongs to>"
    }}
  ],

  "must_study_topics": ["<topic 1>", "<topic 2>"],

  "coverage_summary": "<2-3 sentence summary of what topics dominate>",

  "difficulty_trend": "<one line: is the paper getting harder/easier/stable>"
}}

For topic_frequency: list at least 8-15 topics, sorted by count descending.
For recurring_questions: top 5-10 recurring patterns.
For must_study_topics: top 5-8 topics ranked by importance.

Important:
- Base ALL analysis strictly on the provided papers.
- Do NOT invent statistics — count carefully.
- Output ONLY the JSON. No markdown fences, no explanation.
"""


def build_practice_paper_prompt(
    subject: str, exam: str, analysis_json: str,
    lecture_syllabus: str, num_papers: int,
) -> str:
    t3_special = ""
    if exam == "T3":
        t3_special = """

CRITICAL FOR T3 (End Sem):
- T3 papers cover the FULL syllabus.
- Include ~30% questions from T1 topics (earlier material)
- Include ~30% questions from T2 topics (middle material)
- Include ~40% questions from post-T2 / new topics
- Match this weightage carefully.
"""

    return f"""You are an expert examination paper setter for the subject: **{subject}**
Exam type: **{exam}**

You have access to:
1. Analysis of {num_papers} previous year question papers (JSON below)
2. Lecture syllabus overview
{t3_special}

ANALYSIS JSON:
{analysis_json}

LECTURE SYLLABUS:
{lecture_syllabus}

TASK: Generate a PREDICTED/PRACTICE question paper for the upcoming {exam} exam.

Requirements:
- Match the marks distribution and question type distribution from PYQs.
- Prioritize "must_study_topics" and high-frequency topics.
- Generate NEW questions inspired by recurring PYQ patterns (not exact copies).
- Include a mix of question types (Long, Short, MCQ, Numerical).
- Total marks should match typical PYQ marks.
- Number every question and mention marks in [brackets].
- For each question, provide **answer + brief explanation** IMMEDIATELY below.

FORMAT (strict):

### Q1. [<marks> marks] [<TopicName>]
<question text>

**Answer:** <answer>
**Explanation:** <brief 2-3 line explanation>

---

### Q2. [<marks> marks] [<TopicName>]
...

Continue until total marks reach typical value.
"""


# ============== ANALYSIS EXECUTION ==============

def run_analysis(content: str, subject: str, exam: str, num_papers: int) -> Dict:
    """Send papers to LLM, get structured JSON analysis."""
    MAX_WORDS = 14000
    word_count = len(content.split())
    if word_count > MAX_WORDS:
        words = content.split()
        content = " ".join(words[:MAX_WORDS])
        print(f"   Truncated PYQ content to {MAX_WORDS} words")

    system_prompt = build_analysis_prompt(subject, exam, num_papers)
    user_msg = f"Previous year question papers content:\n\n{content}"

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_msg},
    ]

    print("   Analyzing papers with LLM...")
    raw_response = get_answer(messages, max_tokens=3000)
    raw_response = raw_response.strip()

    if raw_response.startswith("```"):
        raw_response = re.sub(r"^```(?:json)?\s*", "", raw_response)
        raw_response = re.sub(r"\s*```$", "", raw_response)

    match = re.search(r"\{.*\}", raw_response, re.DOTALL)
    if match:
        raw_response = match.group(0)

    try:
        return json.loads(raw_response)
    except json.JSONDecodeError as e:
        print(f"   [WARN] LLM returned non-JSON output: {e}")
        print(f"   Raw output preview: {raw_response[:300]}")
        return {"_parse_error": True, "raw": raw_response}


def generate_practice_paper(
    analysis: Dict, subject: str, exam: str,
    lecture_syllabus: str, num_papers: int,
) -> str:
    """Generate predicted practice paper based on analysis."""
    analysis_str = json.dumps(analysis, indent=2)

    system_prompt = build_practice_paper_prompt(
        subject, exam, analysis_str, lecture_syllabus, num_papers
    )
    user_msg = f"Generate the predicted practice paper for {subject} {exam} now."
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_msg},
    ]

    print("   Generating predicted practice paper...")
    return get_answer(messages, max_tokens=3500)


# ============== INTERACTIVE PARSING ==============

def parse_practice_paper(paper_text: str) -> List[Dict]:
    """Parse practice paper into question dicts."""
    questions = []
    chunks = re.split(r"(?=^#{1,3}\s*Q\d+\.)", paper_text, flags=re.MULTILINE)

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        header_match = re.match(r"#{1,3}\s*(Q\d+)\.?\s*(.*)", chunk)
        if not header_match:
            continue
        qnum = header_match.group(1)
        header = chunk.split("\n", 1)[0].lstrip("#").strip()

        ans_match = re.search(
            r"\*\*Answer:\*\*\s*(.+?)(?=\n\*\*Explanation|\n---|\Z)",
            chunk, flags=re.DOTALL | re.IGNORECASE,
        )
        exp_match = re.search(
            r"\*\*Explanation:\*\*\s*(.+?)(?=\n---|\Z)",
            chunk, flags=re.DOTALL | re.IGNORECASE,
        )
        answer = ans_match.group(1).strip() if ans_match else ""
        explanation = exp_match.group(1).strip() if exp_match else ""

        question_block = chunk
        question_block = re.sub(r"^#{1,3}\s*Q\d+\..*\n", "", question_block)
        question_block = re.sub(
            r"\*\*Answer:\*\*.*?(?=\n\*\*Explanation|\n---|\Z)",
            "", question_block, flags=re.DOTALL | re.IGNORECASE,
        )
        question_block = re.sub(
            r"\*\*Explanation:\*\*.*?(?=\n---|\Z)",
            "", question_block, flags=re.DOTALL | re.IGNORECASE,
        )
        question_block = question_block.replace("---", "").strip()

        questions.append({
            "number": qnum,
            "header": header,
            "question_block": question_block,
            "answer": answer,
            "explanation": explanation,
        })

    return questions