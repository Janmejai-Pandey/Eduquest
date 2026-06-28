import os
import sys
import io
import json
import re
from typing import List, Dict, Optional, Tuple

import requests
import pdfplumber
from pptx import Presentation

from config import chat_imports
chat_imports()
from llm import get_answer

CURRENT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC_DIR = os.path.join(CURRENT_DIR, "..")
# ============== PATHS ==============

BASE_DATASET = os.path.join(CURRENT_DIR, "..", "..", "dataset", "study_material", "CSE")

SEM_FOLDERS = {
    "3": os.path.join(BASE_DATASET, "3"),
    "4": os.path.join(BASE_DATASET, "4"),
}


# ============== QUIZ PROMPT BUILDER ==============

def build_quiz_prompt(
    question_types: List[str],
    difficulty: str,
    num_questions: int,
) -> str:
    """Build a structured system prompt for quiz generation (inline answers)."""
    types_str = ", ".join(question_types)

    return f"""You are an expert quiz-master and academic question setter.

Generate EXACTLY {num_questions} high-quality quiz questions from the
provided content.

REQUIREMENTS:
- Question types to use: {types_str}
- Difficulty level: {difficulty}
- Distribute question types as evenly as possible across the quiz.
- Each question MUST test understanding, not just recall (unless Easy).
- For MCQs: provide exactly 4 options labeled (A), (B), (C), (D).
- For True/False: state a clear factual statement.
- For Fill in the Blanks: use _____ for the blank.
- For Numerical Problems: include all necessary data; show the
  final answer with units.
- For Short Answer: expected answer should fit in 1-2 lines.
- For Long Answer: expected answer should be descriptive (5-10 marks
  style with structured sub-points).

ANSWER POLICY:
Show the correct answer + brief explanation IMMEDIATELY after each
question.

FORMAT (strict):
Use this exact structure for every question:

### Q1. [QuestionType] [Difficulty]
<question text here>

(For MCQ:)
(A) option 1
(B) option 2
(C) option 3
(D) option 4

**Answer:** <answer>
**Explanation:** <brief explanation>

---

Repeat for Q2, Q3, ... up to Q{num_questions}.

Be rigorous, factually accurate, and base questions strictly on the
provided content. Do not invent facts not present in the material.
"""


# ============== LOAD & FILTER ==============

def load_gdrive_links(semester: str) -> List[Dict]:
    sem_folder = SEM_FOLDERS.get(semester)
    if not sem_folder or not os.path.isdir(sem_folder):
        raise FileNotFoundError(f"Semester folder not found: {sem_folder}")
    json_path = os.path.join(sem_folder, "gdrive_links.json")
    if not os.path.isfile(json_path):
        raise FileNotFoundError(f"gdrive_links.json not found at: {json_path}")
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_subjects(links: List[Dict]) -> List[str]:
    return sorted(set(link["subject"] for link in links))


def get_files_by_category(
    links: List[Dict], subject: str, category: str
) -> List[Dict]:
    filtered = [
        link for link in links
        if link["subject"] == subject
        and link["category"] == category
        and not link.get("is_folder", False)
    ]

    def natural_sort_key(item):
        name = item.get("item_name", "")
        numbers = re.findall(r"\d+", name)
        first_num = int(numbers[0]) if numbers else 999
        return (first_num, name.lower())

    filtered.sort(key=natural_sort_key)
    return filtered


def get_available_categories(
    links: List[Dict], subject: str
) -> List[str]:
    cats = set()
    for link in links:
        if link["subject"] == subject and not link.get("is_folder", False):
            cats.add(link["category"])
    ordered = ["Lectures", "Tutorials", "PYQs", "Course Description"]
    return [c for c in ordered if c in cats]


def get_pyq_years(pyq_files: List[Dict]) -> List[str]:
    years = set()
    for f in pyq_files:
        name = f.get("item_name", "")
        match = re.search(r"(20\d{2})", name)
        if match:
            years.add(match.group(1))
        else:
            years.add("Untagged")
    return sorted(years)


def get_pyq_files_by_year(
    pyq_files: List[Dict], year: str
) -> List[Dict]:
    result = []
    for f in pyq_files:
        name = f.get("item_name", "")
        match = re.search(r"(20\d{2})", name)
        tag = match.group(1) if match else "Untagged"
        if tag == year:
            result.append(f)
    return result


# ============== STREAM FROM GDRIVE ==============

def stream_from_gdrive(file_id: str) -> Optional[io.BytesIO]:
    session = requests.Session()
    session.headers.update({
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36"
        )
    })
    base_url = "https://drive.google.com/uc?export=download"
    try:
        response = session.get(
            base_url, params={"id": file_id}, stream=True, timeout=30
        )
        token = None
        for k, v in response.cookies.items():
            if k.startswith("download_warning"):
                token = v
                break
        if not token:
            match = re.search(r'confirm=([0-9A-Za-z_\-]+)', response.text)
            if match:
                token = match.group(1)
        if token:
            response = session.get(
                base_url,
                params={"id": file_id, "confirm": token},
                stream=True, timeout=30,
            )
        buffer = io.BytesIO()
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                buffer.write(chunk)
        buffer.seek(0)
        return buffer
    except Exception as e:
        print(f"   Stream failed: {e}")
        return None


def extract_pdf_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    records = []
    try:
        with pdfplumber.open(buffer) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    records.append({
                        "source_file": source_name,
                        "location": f"page {i}",
                        "text": text,
                    })
    except Exception as e:
        print(f"   PDF extraction failed: {e}")
    return records


def extract_pptx_from_buffer(buffer: io.BytesIO, source_name: str) -> List[Dict]:
    records = []
    try:
        prs = Presentation(buffer)
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            text = "\n".join(parts).strip()
            if text:
                records.append({
                    "source_file": source_name,
                    "location": f"slide {i}",
                    "text": text,
                })
    except Exception as e:
        print(f"   PPTX extraction failed: {e}")
    return records


def extract_text_from_gdrive(file_info: Dict) -> Tuple[str, str]:
    file_id = file_info.get("file_id", "")
    item_name = file_info.get("item_name", "Unknown")
    view_url = file_info.get(
        "original_url",
        f"https://drive.google.com/file/d/{file_id}/view"
    )
    if not file_id:
        return "", view_url

    print(f"   Streaming: {item_name}")
    buffer = stream_from_gdrive(file_id)
    if not buffer:
        return "", view_url

    lower_name = item_name.lower()
    records = []
    if lower_name.endswith(".pdf"):
        records = extract_pdf_from_buffer(buffer, item_name)
    elif lower_name.endswith(".pptx"):
        records = extract_pptx_from_buffer(buffer, item_name)
    else:
        records = extract_pdf_from_buffer(buffer, item_name)
        if not records:
            buffer.seek(0)
            records = extract_pptx_from_buffer(buffer, item_name)
    if not records:
        return "", view_url
    parts = [f"--- {r['location']} ---\n{r['text']}" for r in records]
    return "\n\n".join(parts), view_url


def gather_content(files: List[Dict]) -> Tuple[str, List[Dict]]:
    all_texts = []
    processed_links = []
    for f in files:
        text, view_url = extract_text_from_gdrive(f)
        if not text.strip():
            print(f"   Skipped (no text): {f.get('item_name')}")
            continue
        all_texts.append(
            f"\n{'='*60}\nFILE: {f.get('item_name')}\n{'='*60}\n{text}"
        )
        processed_links.append({
            "name": f.get("item_name", "Unknown"),
            "url": view_url,
        })
    return "\n\n".join(all_texts), processed_links


# ============== AUTO-DECIDE NUM QUESTIONS ==============

def auto_num_questions(text: str) -> int:
    """Decide question count based on content length."""
    words = len(text.split())
    if words < 500:
        return 5
    elif words < 1500:
        return 8
    elif words < 4000:
        return 12
    elif words < 8000:
        return 18
    else:
        return 25


# ============== GENERATE QUIZ ==============

def generate_quiz(
    content: str,
    subject: str,
    label: str,
    question_types: List[str],
    difficulty: str,
    num_questions: int,
) -> str:
    """Call LLM and return the quiz text (inline answer format)."""
    MAX_WORDS = 14000
    word_count = len(content.split())
    if word_count > MAX_WORDS:
        words = content.split()
        content = " ".join(words[:MAX_WORDS])
        print(f"   Truncated content to {MAX_WORDS} words for LLM context")

    system_prompt = build_quiz_prompt(
        question_types, difficulty, num_questions
    )
    user_msg = (
        f"Subject: {subject}\n"
        f"Source: {label}\n\n"
        f"Content to generate quiz from:\n{content}"
    )
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_msg},
    ]
    print("   Generating quiz...")
    return get_answer(messages)


# ============== PARSE QUIZ FOR INTERACTIVE MODE ==============

def parse_quiz(quiz_text: str) -> List[Dict]:
    """Parse quiz text into a list of question dicts."""
    questions = []
    chunks = re.split(r"(?=^#{1,3}\s*Q\d+\.)", quiz_text, flags=re.MULTILINE)

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        header_match = re.match(r"#{1,3}\s*(Q\d+)\.?\s*(.*)", chunk)
        if not header_match:
            continue
        qnum = header_match.group(1)
        header = chunk.split("\n", 1)[0].lstrip("#").strip()

        ans_match = re.search(
            r"\*\*Answer:\*\*\s*(.+?)(?=\n\*\*Explanation|\n---|\Z)",
            chunk, flags=re.DOTALL | re.IGNORECASE,
        )
        exp_match = re.search(
            r"\*\*Explanation:\*\*\s*(.+?)(?=\n---|\Z)",
            chunk, flags=re.DOTALL | re.IGNORECASE,
        )
        answer = ans_match.group(1).strip() if ans_match else ""
        explanation = exp_match.group(1).strip() if exp_match else ""

        question_block = chunk
        question_block = re.sub(r"^#{1,3}\s*Q\d+\..*\n", "", question_block)
        question_block = re.sub(
            r"\*\*Answer:\*\*.*?(?=\n\*\*Explanation|\n---|\Z)",
            "", question_block, flags=re.DOTALL | re.IGNORECASE,
        )
        question_block = re.sub(
            r"\*\*Explanation:\*\*.*?(?=\n---|\Z)",
            "", question_block, flags=re.DOTALL | re.IGNORECASE,
        )
        question_block = question_block.replace("---", "").strip()

        questions.append({
            "number": qnum,
            "header": header,
            "question_block": question_block,
            "answer": answer,
            "explanation": explanation,
        })

    return questions


# ============== SAVE QUIZ ==============

def save_quiz(
    quiz_text: str,
    semester: str,
    subject: str,
    label: str,
    links: Optional[List[Dict]] = None,
) -> str:
    """Save quiz markdown (inline answers always). Returns file path."""
    safe_subject = re.sub(r'[<>:"/\\|?*]', "_", subject).strip()
    out_dir = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "..", "..", "quizzes",
        f"sem_{semester}",
        safe_subject,
    )
    os.makedirs(out_dir, exist_ok=True)

    clean_label = re.sub(r'[<>:"/\\|?*]', "_", label).strip()
    if len(clean_label) > 100:
        clean_label = clean_label[:100]
    filename = f"quiz_{clean_label}.md"
    filepath = os.path.join(out_dir, filename)

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(f"# Quiz: {label}\n")
        f.write(f"**Subject**: {subject}\n")
        f.write(f"**Semester**: {semester}\n\n")
        if links:
            f.write("## Source Files (click to open)\n\n")
            for link in links:
                f.write(f"- [{link['name']}]({link['url']})\n")
            f.write("\n")
        f.write("---\n\n")
        f.write(quiz_text)

    return filepath